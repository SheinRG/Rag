"""
Nexus — Document Ingestion Pipeline
Complete pipeline: download → parse → chunk → embed → store.
Memory-optimised for Render free tier (512 MB RAM).
"""

import os
import gc
import logging
import tempfile
from typing import List

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from database import supabase, embedder
from config import STORAGE_BUCKET, CHUNK_SIZE, CHUNK_OVERLAP

logger = logging.getLogger(__name__)

# Initialize text splitter once
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    separators=["\n\n", "\n", ".", " "],
    length_function=len,
)


def download_from_storage(storage_path: str) -> str:
    """Download a file from Supabase Storage to a temporary local path."""
    try:
        response = supabase.storage.from_(STORAGE_BUCKET).download(storage_path)

        ext = os.path.splitext(storage_path)[1]
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        tmp.write(response)
        tmp.close()

        logger.info(f"Downloaded {storage_path} to {tmp.name}")
        return tmp.name

    except Exception as e:
        logger.error(f"Failed to download from storage: {e}")
        raise


def parse_document(file_path: str, file_type: str) -> List[Document]:
    """Parse a document file into LangChain Document objects."""
    try:
        if file_type == "pdf":
            import fitz
            doc = fitz.open(file_path)
            documents = []
            for i, page in enumerate(doc):
                text = page.get_text() or ""
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={"page": i + 1, "source": file_path}
                    ))
            doc.close()
            if not documents:
                raise ValueError("PDF contains no readable text.")
            return documents

        elif file_type == "docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            text = "\n".join(full_text)
            if not text.strip():
                raise ValueError("Word document is empty.")
            return [Document(page_content=text, metadata={"source": file_path})]

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            documents = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text = "\n".join(slide_text)
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={"page": i + 1, "source": file_path}
                    ))
            if not documents:
                raise ValueError("PowerPoint contains no readable text.")
            return documents

        elif file_type in ("xlsx", "xls"):
            import pandas as pd
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            documents = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                content = df.to_string(index=False)
                if content.strip():
                    documents.append(Document(
                        page_content=f"Sheet: {sheet_name}\n\n{content}",
                        metadata={"sheet": sheet_name, "source": file_path}
                    ))
            if not documents:
                raise ValueError("Excel file is empty.")
            return documents

        elif file_type in ("txt", "md"):
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            if not text.strip():
                raise ValueError("File is empty.")
            return [Document(page_content=text, metadata={"source": file_path})]

        elif file_type == "csv":
            import csv
            documents = []
            with open(file_path, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                for row in reader:
                    content = " | ".join([f"{k}: {v}" for k, v in row.items()])
                    documents.append(Document(page_content=content, metadata={"source": file_path}))
            if not documents:
                raise ValueError("CSV is empty or invalid.")
            return documents

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    except Exception as e:
        logger.error(f"Failed to parse document: {e}")
        raise


def chunk_documents(documents: List[Document]) -> List[Document]:
    """Split documents into smaller chunks."""
    chunks = text_splitter.split_documents(documents)
    logger.info(f"Split into {len(chunks)} chunks.")
    return chunks


def run_ingestion(
    document_id: str,
    storage_path: str,
    file_type: str,
    user_id: str,
):
    """
    Orchestrates the full ingestion pipeline.
    Ultra memory-safe: for large PDFs, processes 50 pages at a time.
    Each page-group is parsed → chunked → embedded → stored → freed.
    This keeps RAM under 512 MB even for 700+ page books.
    """
    local_path = None

    try:
        logger.info(f"Starting ingestion for document {document_id}")

        # Step 1: Download from storage
        local_path = download_from_storage(storage_path)

        total_stored = 0
        chunk_index = 0

        if file_type == "pdf":
            # ── STREAMING PDF INGESTION (page-group at a time) ──
            import fitz
            pdf_doc = fitz.open(local_path)
            total_pages = len(pdf_doc)
            pdf_doc.close()
            del pdf_doc
            gc.collect()

            logger.info(f"PDF has {total_pages} pages. Processing in groups of 50.")

            PAGE_GROUP = 50
            for page_start in range(0, total_pages, PAGE_GROUP):
                page_end = min(page_start + PAGE_GROUP, total_pages)

                # Parse only this group of pages
                pdf_doc = fitz.open(local_path)
                documents = []
                for i in range(page_start, page_end):
                    page = pdf_doc[i]
                    text = page.get_text() or ""
                    if text.strip():
                        documents.append(Document(
                            page_content=text,
                            metadata={"page": i + 1, "source": local_path}
                        ))
                pdf_doc.close()
                del pdf_doc
                gc.collect()

                if not documents:
                    continue

                # Chunk this group
                chunks = text_splitter.split_documents(documents)
                del documents
                gc.collect()

                if not chunks:
                    continue

                # Embed and store
                texts = [c.page_content for c in chunks]
                embeddings = [e.tolist() for e in embedder.embed(texts)]

                rows = []
                for c, emb in zip(chunks, embeddings):
                    rows.append({
                        "document_id": document_id,
                        "user_id": user_id,
                        "content": c.page_content,
                        "embedding": emb,
                        "chunk_index": chunk_index,
                    })
                    chunk_index += 1

                supabase.table("chunks").insert(rows).execute()
                total_stored += len(rows)

                logger.info(
                    f"Pages {page_start+1}-{page_end}: "
                    f"{len(rows)} chunks stored ({total_stored} total)"
                )

                del chunks, texts, embeddings, rows
                gc.collect()

        else:
            # ── NON-PDF FILES (small, process normally) ──
            documents = parse_document(local_path, file_type)

            if not documents or all(not d.page_content.strip() for d in documents):
                raise ValueError("Document contains no readable text content.")

            chunks = chunk_documents(documents)
            del documents
            gc.collect()

            if not chunks:
                raise ValueError("Document produced no text chunks after splitting.")

            texts = [c.page_content for c in chunks]
            embeddings = [e.tolist() for e in embedder.embed(texts)]

            rows = []
            for i, (c, emb) in enumerate(zip(chunks, embeddings)):
                rows.append({
                    "document_id": document_id,
                    "user_id": user_id,
                    "content": c.page_content,
                    "embedding": emb,
                    "chunk_index": i,
                })

            # Insert in batches of 100
            for b in range(0, len(rows), 100):
                supabase.table("chunks").insert(rows[b:b+100]).execute()

            total_stored = len(rows)
            del chunks, texts, embeddings, rows
            gc.collect()

        if total_stored == 0:
            raise ValueError("Document contains no readable text content.")

        # Update document status to ready
        supabase.table("documents").update({
            "status": "ready",
            "num_chunks": total_stored,
        }).eq("id", document_id).execute()

        logger.info(f"Ingestion complete for document {document_id}: {total_stored} chunks.")

    except Exception as e:
        logger.error(f"Ingestion failed for document {document_id}: {e}")

        try:
            supabase.table("documents").update({
                "status": "failed",
                "error_msg": str(e)[:500],
            }).eq("id", document_id).execute()
        except Exception as db_err:
            logger.error(f"Failed to update document status: {db_err}")

    finally:
        if local_path and os.path.exists(local_path):
            try:
                os.unlink(local_path)
            except Exception:
                pass
